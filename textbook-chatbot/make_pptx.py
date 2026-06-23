#!/usr/bin/env python3
"""
PowerPoint生成ツール
使い方:
  python3 ~/textbook-chatbot/make_pptx.py input.json
  python3 ~/textbook-chatbot/make_pptx.py input.json --out ~/Desktop/output.pptx

JSON形式:
{
  "title": "プレゼンタイトル",
  "subtitle": "サブタイトル（任意）",
  "slides": [
    {
      "title": "スライドタイトル",
      "bullets": ["箇条書き1", "箇条書き2"],
      "body": "本文テキスト（bulletsの代わりに）",
      "note": "発表者ノート（任意）"
    }
  ],
  "references": ["参考文献1", "参考文献2"]
}
"""

import json, sys, argparse, re
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = Path.home() / "Desktop"

# ── カラー定義 ──────────────────────────────────────────────────
C_BLUE_DARK  = RGBColor(0x1d, 0x4e, 0xd8)
C_BLUE_MID   = RGBColor(0x37, 0x6b, 0xe0)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK       = RGBColor(0x1e, 0x29, 0x3b)
C_GRAY       = RGBColor(0x44, 0x55, 0x66)
C_LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xff)

# スライドサイズ（ワイド 16:9）
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_font(run, size: int, bold: bool = False, color: RGBColor = None):
    run.font.size  = Pt(size)
    run.font.bold  = bold
    if color:
        run.font.color.rgb = color


def add_rect(slide, l, t, w, h, color: RGBColor):
    """塗りつぶし矩形を追加"""
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text: str, size: int,
                bold=False, color: RGBColor = C_DARK,
                align=PP_ALIGN.LEFT, wrap=True) -> object:
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p     = tf.paragraphs[0]
    p.alignment = align
    run   = p.add_run()
    run.text = text
    set_font(run, size, bold, color)
    return txBox


def make_title_slide(prs: Presentation, title: str, subtitle: str, date_str: str):
    """タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 背景（濃い青）
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BLUE_DARK)

    # アクセントライン（細い白線）
    line = slide.shapes.add_shape(1,
        Inches(1.0), Inches(3.8), Inches(11.33), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_WHITE
    line.line.fill.background()

    # タイトル
    add_textbox(slide,
        Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.8),
        title, 36, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # サブタイトル
    if subtitle:
        add_textbox(slide,
            Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.0),
            subtitle, 20, color=C_LIGHT_BLUE, align=PP_ALIGN.LEFT)

    # 日付
    add_textbox(slide,
        Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
        date_str, 14, color=C_LIGHT_BLUE, align=PP_ALIGN.RIGHT)


def make_content_slide(prs: Presentation, slide_data: dict, slide_num: int, total: int):
    """コンテンツスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── ヘッダーバー ───────────────────────────────────────────
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_BLUE_DARK)

    # スライドタイトル
    add_textbox(slide,
        Inches(0.4), Inches(0.15), Inches(11.5), Inches(0.8),
        slide_data.get("title", ""), 24, bold=True, color=C_WHITE)

    # ページ番号
    add_textbox(slide,
        Inches(11.8), Inches(0.2), Inches(1.3), Inches(0.6),
        f"{slide_num}/{total}", 12, color=C_LIGHT_BLUE, align=PP_ALIGN.RIGHT)

    # ── コンテンツエリア ───────────────────────────────────────
    content_top  = Inches(1.25)
    content_left = Inches(0.5)
    content_w    = Inches(12.33)
    content_h    = Inches(5.5)

    bullets = slide_data.get("bullets", [])
    body    = slide_data.get("body", "")

    if bullets:
        # 箇条書きモード
        txBox = slide.shapes.add_textbox(content_left, content_top, content_w, content_h)
        tf    = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # インデントレベル判定（先頭の "  " や "　" でレベル分け）
            level = 0
            text  = bullet
            if bullet.startswith("  ") or bullet.startswith("　"):
                level = 1
                text  = bullet.lstrip()

            p.level = level
            run     = p.add_run()
            run.text = ("• " if level == 0 else "  ‐ ") + text
            size    = 20 if level == 0 else 17
            color   = C_DARK if level == 0 else C_GRAY
            set_font(run, size, bold=(level == 0), color=color)
            p.space_before = Pt(6 if level == 0 else 2)

    elif body:
        # 本文モード
        add_textbox(slide, content_left, content_top, content_w, content_h,
                    body, 18, color=C_DARK)

    # ── 発表者ノート ───────────────────────────────────────────
    note_text = slide_data.get("note", "")
    if note_text:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = note_text

    # ── フッター（出典）───────────────────────────────────────
    source = slide_data.get("source", "")
    if source:
        add_textbox(slide,
            Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.4),
            f"出典: {source}", 9, color=C_GRAY)


def make_reference_slide(prs: Presentation, refs: list):
    """参考文献スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_BLUE_DARK)
    add_textbox(slide, Inches(0.4), Inches(0.15), Inches(11.5), Inches(0.8),
                "参考文献", 24, bold=True, color=C_WHITE)

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, ref in enumerate(refs):
        p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"・{ref}"
        set_font(run, 14, color=C_DARK)
        p.space_before = Pt(6)


def generate_pptx(data: dict, out_path: Path):
    """PowerPointを生成"""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    title    = data.get("title", "無題")
    subtitle = data.get("subtitle", "")
    slides   = data.get("slides", [])
    refs     = data.get("references", [])
    date_str = datetime.now().strftime("%Y年%m月%d日")

    # タイトルスライド
    make_title_slide(prs, title, subtitle, date_str)

    # コンテンツスライド
    total = len(slides)
    for i, s in enumerate(slides, 1):
        make_content_slide(prs, s, i, total)

    # 参考文献スライド
    if refs:
        make_reference_slide(prs, refs)

    prs.save(str(out_path))
    return out_path


def main():
    parser = argparse.ArgumentParser(description="PowerPoint生成")
    parser.add_argument("input", help="入力JSONファイル（または - でstdin）")
    parser.add_argument("--out", "-o", default="", help="出力ファイルパス")
    args = parser.parse_args()

    if args.input == "-":
        data = json.load(sys.stdin)
    else:
        with open(args.input, encoding="utf-8") as f:
            data = json.load(f)

    if args.out:
        out_path = Path(args.out)
    else:
        ts   = datetime.now().strftime("%Y%m%d_%H%M")
        name = re.sub(r'[^\w぀-鿿]', '_', data.get("title", "output"))
        out_path = OUTPUT_DIR / f"{name}_{ts}.pptx"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    result = generate_pptx(data, out_path)
    print(f"✅ PowerPoint保存完了: {result}")


if __name__ == "__main__":
    main()
